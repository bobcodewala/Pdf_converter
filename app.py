from flask import Flask, request, send_file
from PIL import Image
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from docx import Document
import os

app = Flask(__name__)

# Configure the upload folder
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


def convert_txt_to_pdf(text_content, c):
    """Convert .txt file content to PDF."""
    lines = text_content.splitlines()
    y_position = 750
    for line in lines:
        c.drawString(100, y_position, line)
        y_position -= 15
        if y_position <= 40:
            c.showPage()
            y_position = 750
    c.showPage()


def convert_docx_to_pdf(file_path, c):
    """Convert .docx to PDF."""
    doc = Document(file_path)
    y_position = 750
    for para in doc.paragraphs:
        c.drawString(100, y_position, para.text)
        y_position -= 15
        if y_position <= 40:
            c.showPage()
            y_position = 750
    c.showPage()


def convert_pptx_to_pdf(file_path, c):
    """Convert .pptx to PDF."""
    ppt = Presentation(file_path)
    y_position = 750
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    c.drawString(100, y_position, para.text)
                    y_position -= 15
                    if y_position <= 40:
                        c.showPage()
                        y_position = 750
        c.showPage()


def convert_image_to_pdf(file_path, c):
    """Convert image to PDF."""
    img = Image.open(file_path)
    img_width, img_height = img.size
    img_width = min(img_width, 500)  # Resize to fit within the PDF page
    img_height = min(img_height, 700)
    c.drawImage(file_path, 100, 500, width=img_width, height=img_height)
    c.showPage()


@app.route('/')
def index():
    return app.send_static_file('index.html')


@app.route('/convert-to-pdf', methods=['POST'])
def convert_to_pdf():
    files = request.files.getlist('files')
    pdf_stream = BytesIO()
    c = canvas.Canvas(pdf_stream, pagesize=letter)

    # Process each file
    for file in files:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)

        if file.filename.lower().endswith('.txt'):
            with open(file_path, 'r') as f:
                text_content = f.read()
            convert_txt_to_pdf(text_content, c)
        elif file.filename.lower().endswith('.docx'):
            convert_docx_to_pdf(file_path, c)
        elif file.filename.lower().endswith('.pptx'):
            convert_pptx_to_pdf(file_path, c)
        elif file.filename.lower().endswith(('jpg', 'jpeg', 'png', 'gif', 'bmp')):
            convert_image_to_pdf(file_path, c)

    # Save the PDF
    c.save()
    pdf_stream.seek(0)

    return send_file(pdf_stream, as_attachment=True, download_name="converted_files.pdf", mimetype="application/pdf")


if __name__ == "__main__":
    app.run(debug=True)
